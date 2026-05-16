from dotenv import load_dotenv
load_dotenv()
import os
import re
from pathlib import Path
from typing import List, Dict, Any
from datetime import datetime



class SimpleDocumentIngestor:
    """Simple document ingestor that doesn't require heavy dependencies"""
    
    def __init__(self, raw_docs_folder: str = "./raw_docs"):
        self.raw_docs_folder = Path(raw_docs_folder)
        self.processed_files = self._load_processed_files()
    
    def _load_processed_files(self) -> set:
        """Load list of already processed files"""
        processed_file = self.raw_docs_folder / ".processed_files.txt"
        if processed_file.exists():
            with open(processed_file, 'r') as f:
                return set(line.strip() for line in f)
        return set()
    
    def _save_processed_file(self, filename: str):
        """Mark a file as processed"""
        processed_file = self.raw_docs_folder / ".processed_files.txt"
        with open(processed_file, 'a') as f:
            f.write(filename + "\n")
        self.processed_files.add(filename)
    
    def extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file using the knowledge base method"""
        try:
            from knowledge.vector_store import get_knowledge_base
            kb = get_knowledge_base()
            # This will add the PDF directly with proper parsing
            kb.add_pdf_document(str(file_path))
            return ""  # Return empty since it's handled by add_pdf_document
        except Exception as e:
            print(f"   Error reading PDF: {e}")
            return ""
    
    def extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from Word document"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except ImportError:
            print("   ⚠️ python-docx not installed. Run: pip install python-docx")
            return ""
        except Exception as e:
            print(f"   Error reading DOCX: {e}")
            return ""
        
    
    
    def extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            print("   ⚠️ python-pptx not installed. Run: pip install python-pptx")
            return ""
        except Exception as e:
            print(f"   Error reading PPTX: {e}")
            return ""
    
    def extract_text_from_txt(self, file_path: Path) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                print(f"   Error reading TXT: {e}")
                return ""
            
    
    def extract_text_from_excel(self, file_path: Path) -> str:
        """Extract text from Excel file (.xlsx, .xls)"""
        try:
            import pandas as pd
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            all_text = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                all_text.append(f"\n=== Sheet: {sheet_name} ===\n")
                # Convert dataframe to readable text
                all_text.append(df.to_string(index=False))
            
            return "\n".join(all_text)
        except ImportError:
            print("   ⚠️ pandas not installed. Run: pip install pandas openpyxl")
            return ""
        except Exception as e:
            print(f"   Error reading Excel: {e}")
            return ""

    def extract_text_from_csv(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except ImportError:
            print("   ⚠️ pandas not installed. Run: pip install pandas")
            return ""
        except Exception as e:
            print(f"   Error reading CSV: {e}")
            return ""
    
    def detect_category(self, filename: str, content: str = "") -> str:
        """Auto-detect category based on filename and content"""
        filename_lower = filename.lower()
        content_lower = content.lower()
        
        # Check filename patterns
        if any(word in filename_lower for word in ['program', 'course', 'degree', 'bachelor', 'master', 'phd']):
            return "programs"
        elif any(word in filename_lower for word in ['fee', 'tuition', 'cost', 'payment', 'bursary']):
            return "fees"
        elif any(word in filename_lower for word in ['admission', 'apply', 'entry', 'requirement', 'application']):
            return "admissions"
        elif any(word in filename_lower for word in ['regulation', 'policy', 'rule', 'guideline', 'handbook']):
            return "regulations"
        elif any(word in filename_lower for word in ['calendar', 'schedule', 'date', 'term']):
            return "calendar"
        elif any(word in filename_lower for word in ['scholarship', 'grant', 'sponsor']):
            return "scholarships"
        
        # Check content if available
        if content_lower:
            if any(word in content_lower for word in ['programme', 'course of study', 'degree programme']):
                return "programs"
            elif any(word in content_lower for word in ['fee structure', 'tuition fees', 'how to pay']):
                return "fees"
        
        return "general"
    
    def process_file(self, file_path: Path) -> Dict[str, Any]:
        """Process a single file and extract its content"""
        file_ext = file_path.suffix.lower()
        
        print(f"   Extracting text from {file_ext} file...")
        
        if file_ext == '.pdf':
            content = self.extract_text_from_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            content = self.extract_text_from_docx(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            content = self.extract_text_from_pptx(file_path)
        elif file_ext == '.txt':
            content = self.extract_text_from_txt(file_path)
        elif file_ext in ['.xlsx', '.xls']:  # Add Excel support
            content = self.extract_text_from_excel(file_path)
        elif file_ext == '.csv':  # Add CSV support
            content = self.extract_text_from_csv(file_path)
        else:
            return None
        
        if not content or len(content.strip()) < 50:
            print(f"   ⚠️ No text extracted (content too short or empty)")
            return None
        
        # Detect category
        category = self.detect_category(file_path.name, content)
        
        # Generate title
        title = file_path.stem.replace('_', ' ').replace('-', ' ').title()
        
        # Limit content length to avoid huge vectors
        if len(content) > 10000:
            content = content[:10000]
            print(f"   📝 Content truncated to 10000 characters")
        
        return {
            "title": title,
            "category": category,
            "content": content,
            "source_file": file_path.name
        }
    
    def ingest_all(self) -> Dict[str, Any]:
        """Ingest all documents from the folder"""
        if not self.raw_docs_folder.exists():
            print(f"📁 Creating folder: {self.raw_docs_folder}")
            self.raw_docs_folder.mkdir(parents=True)
            print(f"   Please place your documents in: {self.raw_docs_folder.absolute()}")
            return {"processed": 0, "total_found": 0, "documents": []}
        
        # Find all supported documents
        supported_extensions = ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.xlsx', '.xls', '.csv']
        documents = []
        
        for ext in supported_extensions:
            documents.extend(self.raw_docs_folder.glob(f"*{ext}"))
        
        if not documents:
            print(f"📂 No documents found in: {self.raw_docs_folder}")
            print(f"   Supported formats: PDF, DOCX, PPTX, TXT")
            return {"processed": 0, "total_found": 0, "documents": []}
        
        print(f"\n📄 Found {len(documents)} document(s)")
        
        from knowledge.vector_store import get_knowledge_base
        kb = get_knowledge_base()
        
        new_documents = []
        skipped = []
        
        for doc_path in documents:
            if doc_path.name in self.processed_files:
                skipped.append(doc_path.name)
                continue
            
            print(f"\n📖 Processing: {doc_path.name}")
            doc_data = self.process_file(doc_path)
            
            if doc_data:
                kb.add_document(
                    title=doc_data['title'],
                    content=doc_data['content'],
                    category=doc_data['category'],
                    metadata={"source_file": doc_data['source_file']}
                )
                new_documents.append(doc_data)
                self._save_processed_file(doc_path.name)
                print(f"   ✅ Added to knowledge base (Category: {doc_data['category']})")
        
        return {
            "processed": len(new_documents),
            "skipped": len(skipped),
            "total_found": len(documents),
            "documents": [d['title'] for d in new_documents]
        }

# Singleton
_ingestor = None

def get_ingestor():
    global _ingestor
    if _ingestor is None:
        _ingestor = SimpleDocumentIngestor()
    return _ingestor